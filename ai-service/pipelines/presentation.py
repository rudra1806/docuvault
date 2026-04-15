# ============================================================
# pipelines/presentation.py — Presentation Extraction Pipeline
# ============================================================
# Extracts text from PPTX slides including titles, body,
# shapes, and speaker notes.
# ============================================================

import io
import logging
from pptx import Presentation

logger = logging.getLogger(__name__)


async def extract_presentation(file_bytes: bytes, file_name: str, ext: str) -> dict:
    """Extract text from PowerPoint presentations."""

    if ext == "ppt":
        # Legacy PPT format — limited support
        return {
            "text": "",
            "source": "presentation",
            "metadata": {
                "type": "ppt",
                "note": "Legacy PPT format has limited extraction support. Please convert to PPTX.",
            },
        }

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        slide_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []

            # Extract title if present
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_text_parts.append(f"Title: {slide.shapes.title.text.strip()}")

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Skip the title since we already got it
                    if shape == slide.shapes.title:
                        continue
                    slide_text_parts.append(shape.text.strip())

                # Extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_data = [
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        ]
                        if row_data:
                            slide_text_parts.append(" | ".join(row_data))

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text_parts.append(f"[Speaker Notes]: {notes}")

            if slide_text_parts:
                text_parts.append(
                    f"[Slide {slide_num}]\n" + "\n".join(slide_text_parts)
                )
                slide_count += 1

        return {
            "text": "\n\n".join(text_parts),
            "source": "presentation",
            "metadata": {
                "slides": slide_count,
                "total_slides": len(prs.slides),
                "type": ext,
            },
        }

    except Exception as e:
        logger.error(f"Presentation extraction failed for {file_name}: {e}")
        return {
            "text": "",
            "source": "presentation",
            "metadata": {"error": str(e), "type": ext},
        }
